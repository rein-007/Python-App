import os
from datetime import datetime
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def build_pptx(kpis, rev_exp_fig, cashflow_fig):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Financial KPI Report"
    slide.placeholders[1].text = "Generated from Streamlit dashboard"

    # KPI slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Key KPIs"
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.0), Inches(3.0))
    tf = tb.text_frame
    tf.clear()
    for key, val in kpis.items():
        p = tf.add_paragraph()
        p.text = f"{key.capitalize()}: {val:.2f}" if val is not None else f"{key.capitalize()}: --"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0,0,0)

    # Chart slides
    for fig, title in [(rev_exp_fig, "Revenue vs Expenses"), (cashflow_fig, "Cash Flow")]:
        img_path = "_temp.png"
        fig.savefig(img_path, bbox_inches="tight")
        plt.close(fig)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(img_path, Inches(0.7), Inches(1.5), Inches(8.5), Inches(4.8))
        os.remove(img_path)

    return prs

def save_pptx(prs, folder="."):
    ts = datetime.now().strftime("%Y-%m-%d_%H-%M")
    filename = f"Financial_KPI_Report_{ts}.pptx"
    path = os.path.join(folder, filename)
    prs.save(path)
    return path