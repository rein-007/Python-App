import os
from datetime import datetime
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from config import ACCENT_COLOR, FONT_NAME

def export_to_pptx(ticker, hist, price, pct_change, folder="."):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"{ticker} Stock Report"

    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(2))
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Price: {price:.2f} USD\nChange: {pct_change:.2f}%"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0,0,0)

    fig, ax = plt.subplots(figsize=(6,4))
    hist['Close'].plot(ax=ax, color=ACCENT_COLOR)
    ax.set_title(f"{ticker} Closing Prices (Last 30 Days)", fontname=FONT_NAME)
    ax.set_ylabel("Price (USD)")
    ax.grid(True)
    fig.savefig("chart.png")

    slide.shapes.add_picture("chart.png", Inches(5), Inches(1.5), Inches(4.5), Inches(3))

    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M")
    filename = f"{ticker}_Stock_Report_{timestamp}.pptx"
    save_path = os.path.join(folder, filename)
    prs.save(save_path)
    return save_path